
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import datetime

# --- Constants & Configuration ---
SLIDE_WIDTH = Inches(13.333) # Widescreen 16:9
SLIDE_HEIGHT = Inches(7.5)
MARGIN_LEFT = Cm(0.5) # Reduced margin to fit more columns
MARGIN_TOP = Cm(2.0)
ROW_HEIGHT = Cm(1.2)
HEADER_HEIGHT = Cm(1.2) # Minimum heights
ROW_HEIGHT_MIN = Cm(1.2)

def estimate_text_lines(text, col_width_cm, font_size_pt):
    """
    Very rough estimation of how many lines text will occupy.
    """
    if not text:
        return 1
    
    # Assume Average character width is roughly 0.6 * font_size in points
    # 1 Pt = 1/72 inch = 0.03527 cm
    avg_char_w_cm = (font_size_pt * 0.5) * 0.03527
    col_width_chars = col_width_cm / avg_char_w_cm
    
    lines_total = 0
    paragraphs = str(text).split('\n')
    for p in paragraphs:
        if not p:
            lines_total += 1
            continue
        # Estimate wrap lines
        wrap_count = (len(p) // int(col_width_chars)) + 1
        lines_total += wrap_count
    return lines_total

# Column Widths 
# Added "Subject" (主題) at index 0
# Info Column Widths (Adjusted for 15-char Task and tighter grid)
INFO_COL_WIDTHS = [
    Cm(3.7),   # 主題 (Subject): 7 chars * 0.45 + 0.5
    Cm(1.8),   # 用戶 (User): 3 chars * 0.45 + 0.5
    Cm(1.8),   # IT窗口 (IT Contact): 3 chars * 0.45 + 0.5
    Cm(2.45),   # 需求單號 (Req ID): 4 chars * 0.45 + 0.5 + 0.15 (手動估算)
    Cm(6.35),  # Task (Task Description): 13 chars * 0.45 + 0.5
    Cm(2.3),   # Status: 4 chars * 0.45 + 0.5
]

# Date Column Widths: 25 cols * 0.58cm = 14.5cm (0.58 已是最小沒法再小了)
# Total Table Width: 18.4cm + 14.5cm = 32.9cm (on 33.87cm slide)
DAY_COL_WIDTH = Cm(0.58)
DATE_COL_WIDTHS = [DAY_COL_WIDTH] * 25

COL_WIDTHS = INFO_COL_WIDTHS + DATE_COL_WIDTHS

HEADERS = ['主題', '用戶', 'IT窗口', '需求單號', 'Task', 'Status']

def get_week_range_str(date_obj):
    """Returns 'W05(01/26-01/30)' format"""
    year, week, weekday = date_obj.isocalendar()
    monday = date_obj - datetime.timedelta(days=weekday-1)
    friday = monday + datetime.timedelta(days=4)
    return f"W{week:02d}\n({monday.strftime('%m/%d')}-{friday.strftime('%m/%d')})" 

def generate_date_headers(base_date):
    """Generates 5 week headers centered on base_date"""
    start_date = base_date - datetime.timedelta(weeks=1)
    headers = []
    current = start_date
    for _ in range(5):
        headers.append(get_week_range_str(current))
        current += datetime.timedelta(weeks=1)
    return headers

def create_pptx(data, today_date=None):
    if today_date is None:
        today_date = datetime.date.today()
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # --- 1. Draw Title ---
    title_shape = slide.shapes.add_textbox(MARGIN_LEFT, Cm(0.5), SLIDE_WIDTH - MARGIN_LEFT*2, Cm(1.5))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('topic', '專案甘特圖')
    p.font.size = Pt(28) # Increased Title
    p.font.bold = True
    
    # --- 2. Process Dates ---
    base_date_str = data.get('base_date', datetime.date.today().strftime('%Y-%m-%d'))
    base_date = datetime.datetime.strptime(base_date_str, '%Y-%m-%d').date()
    today_date = datetime.date.today()
    
    date_headers = generate_date_headers(base_date)
    
    # --- 3. Draw Table ---
    # Rows: Task rows + 1 Header
    # Cols: 6 Info + 25 Dates
    rows_count = len(data.get('tasks', [])) + 1
    cols_count = len(COL_WIDTHS)
    
    total_width = sum(COL_WIDTHS)
    # We use a standard table height and let it expand if needed
    # Header 1.2cm, Data rows 1.2cm+
    table_shape = slide.shapes.add_table(rows_count, cols_count, MARGIN_LEFT, MARGIN_TOP, total_width, Cm(rows_count * 1.2)).table
    
    # Set Column Widths
    for idx, width in enumerate(COL_WIDTHS):
        table_shape.columns[idx].width = width
        
    # Fill Headers
    # Info Headers
    for idx, text in enumerate(HEADERS):
        cell = table_shape.cell(0, idx)
        cell.text = text
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(12)
            p.font.bold = True
    
    # Date Headers (Merged by Week)
    # W05 is col 6-10, W06 is 11-15, etc.
    for i, w_text in enumerate(date_headers):
        start_col = 6 + (i * 5)
        end_col = start_col + 4
        
        # Set text in the first cell of the range
        main_cell = table_shape.cell(0, start_col)
        main_cell.text = w_text
        for p in main_cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(10)
            p.font.bold = True
            
        # Merge if possible (python-pptx merge can be tricky with complex tables, 
        # but for a simple add_table it works well)
        try:
            main_cell.merge(table_shape.cell(0, end_col))
        except:
            pass # Fallback if merge fails

    # Calculate shared date constants for the slide
    # Align to Monday of first week column (Week 1 of the 5-week range)
    first_col_week_date = base_date - datetime.timedelta(weeks=1)
    start_monday = first_col_week_date - datetime.timedelta(days=first_col_week_date.weekday())

    # Fill Data Rows
    for r_idx, task in enumerate(data.get('tasks', [])):
        r = r_idx + 1
        
        # Text Fields
        fields = [
            task.get('subject', ''),
            task.get('user', ''),
            task.get('it_contact', ''),
            task.get('req_id', ''),
            task.get('task_desc', ''),
            task.get('status', '')
        ]
        
        for c_idx, text in enumerate(fields):
            cell = table_shape.cell(r, c_idx)
            cell.vertical_anchor = MSO_SHAPE.RECTANGLE 
            
            if c_idx == 4 and isinstance(text, list):
                tf = cell.text_frame
                tf.text = "" 
                for i, line in enumerate(text):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = str(line)
                    p.font.size = Pt(10)
                    p.level = 0
                    if not p.text.startswith("•"):
                        p.text = "• " + p.text
            else:
                cell.text = str(text) if text is not None else ""
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(10)

        # --- Grid-Based Coloring (The "Gantt Bar") ---
        bar_start_str = task.get('start_date')
        bar_end_str = task.get('end_date')
        bar_text = task.get('bar_text', '')
        
        if bar_start_str and bar_end_str:
            try:
                # Convert dates
                b_start = datetime.datetime.strptime(bar_start_str, '%Y-%m-%d').date()
                b_end = datetime.datetime.strptime(bar_end_str, '%Y-%m-%d').date()

                # Identify grid range (dates to 0-24 index)
                date_indices = []
                for d_idx in range(25):
                    week_num = d_idx // 5
                    day_in_week = d_idx % 5
                    current_week_monday = start_monday + datetime.timedelta(weeks=week_num)
                    actual_date = current_week_monday + datetime.timedelta(days=day_in_week)
                    if b_start <= actual_date <= b_end:
                        date_indices.append((d_idx, actual_date))
                
                if date_indices:
                    past_indices = [idx for idx, dt in date_indices if dt <= today_date]
                    future_indices = [idx for idx, dt in date_indices if dt > today_date]
                    
                    # Merge and color segments
                    for segment in [past_indices, future_indices]:
                        if not segment: continue
                        s_idx, e_idx = segment[0], segment[-1]
                        segment_cell = table_shape.cell(r, 6 + s_idx)
                        if s_idx != e_idx:
                            try: segment_cell.merge(table_shape.cell(r, 6 + e_idx))
                            except: pass
                        
                        segment_cell.fill.solid()
                        if segment == past_indices:
                            segment_cell.fill.fore_color.rgb = RGBColor(91, 155, 213) # Blue
                            if bar_text:
                                segment_cell.text = bar_text
                                for p in segment_cell.text_frame.paragraphs:
                                    p.alignment = PP_ALIGN.CENTER
                                    p.font.size = Pt(9)
                                    p.font.color.rgb = RGBColor(255, 255, 255)
                        else:
                            segment_cell.fill.fore_color.rgb = RGBColor(221, 235, 247) # Light Blue
                            if bar_text and not past_indices:
                                segment_cell.text = bar_text
                                for p in segment_cell.text_frame.paragraphs:
                                    p.alignment = PP_ALIGN.CENTER
                                    p.font.size = Pt(9)
                                    p.font.color.rgb = RGBColor(65, 113, 156)
            except Exception as e:
                print(f"Error drawing bar for task {r_idx}: {e}")
                pass

    return prs

if __name__ == "__main__":
    test_data = {
         'tasks': [
             {'subject': 'S1', 'user': 'U1', 'start_date': '2026-02-02', 'end_date': '2026-02-06', 'bar_text': 'Task 1', 'milestones': '02/03:Check'}
         ]
    }
    create_pptx(test_data).save('d:/localgit/test2.pptx')
